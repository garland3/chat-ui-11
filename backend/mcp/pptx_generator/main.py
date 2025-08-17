"""
PowerPoint Generator MCP Server using FastMCP.

Converts JSON input containing slide data into a PowerPoint presentation.
Input should be a JSON object with a 'slides' array, each containing 'title' and 'content' fields.
Supports bullet point lists in content for more complex slide formatting.

Tools:
 - json_to_pptx: Converts JSON input to PowerPoint presentation

Demonstrates: JSON input handling, file output with base64 encoding, and structured output.
"""

from __future__ import annotations

import base64
import os
import tempfile
import io
import re
import requests
from typing import Any, Dict, List, Annotated, Optional
from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pdfkit import from_string
from PIL import Image


mcp = FastMCP("pptx_generator")


def _is_backend_download_path(s: str) -> bool:
    """Detect backend-relative download paths like /api/files/download/...."""
    return isinstance(s, str) and s.startswith("/api/files/download/")


def _backend_base_url() -> str:
    """Resolve backend base URL from environment variable."""
    return os.environ.get("CHATUI_BACKEND_BASE_URL", "http://127.0.0.1:8000")


def _load_image_bytes(filename: str, file_data_base64: str = "") -> Optional[bytes]:
    """Load image data from filename or base64 data."""
    if file_data_base64:
        try:
            return base64.b64decode(file_data_base64)
        except Exception as e:
            print(f"Error decoding base64 image data: {e}")
            return None
    
    if _is_backend_download_path(filename):
        # Backend provided a download path
        full_url = _backend_base_url() + filename
        try:
            print(f"Fetching image from {full_url}")
            response = requests.get(full_url, timeout=30)
            response.raise_for_status()
            return response.content
        except Exception as e:
            print(f"Error fetching image from {full_url}: {e}")
            return None
    
    # Try as local file path
    if os.path.isfile(filename):
        try:
            with open(filename, "rb") as f:
                return f.read()
        except Exception as e:
            print(f"Error reading local image file {filename}: {e}")
            return None
    
    print(f"Image file not found: {filename}")
    return None


def _parse_markdown_slides(markdown_content: str) -> List[Dict[str, str]]:
    """Parse markdown content into slides."""
    slides = []
    
    # Split by headers (# or ##)
    sections = re.split(r'^#{1,2}\s+(.+)$', markdown_content, flags=re.MULTILINE)
    
    # Remove empty first element if exists
    if sections and not sections[0].strip():
        sections = sections[1:]
    
    # Group into title/content pairs
    for i in range(0, len(sections), 2):
        if i + 1 < len(sections):
            title = sections[i].strip()
            content = sections[i + 1].strip() if i + 1 < len(sections) else ""
            slides.append({"title": title, "content": content})
        elif sections[i].strip():
            # Handle case where there's a title but no content
            slides.append({"title": sections[i].strip(), "content": ""})
    
    # If no headers found, treat entire content as one slide
    if not slides and markdown_content.strip():
        slides.append({"title": "Slide 1", "content": markdown_content.strip()})
    
    return slides


def _add_image_to_slide(slide_obj, image_bytes: bytes, left: Inches = Inches(1), top: Inches = Inches(2), 
                       width: Inches = Inches(8), height: Inches = Inches(5)):
    """Add image to a slide."""
    try:
        # Create a temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_bytes)
            tmp_file.flush()
            
            # Add image to slide
            pic = slide_obj.shapes.add_picture(tmp_file.name, left, top, width, height)
            
            # Clean up
            os.unlink(tmp_file.name)
            
            return pic
    except Exception as e:
        print(f"Error adding image to slide: {e}")
        return None


@mcp.tool
def json_to_pptx(
    input_data: Annotated[str, "JSON string containing slide data in this format: {\"slides\": [{\"title\": \"Slide 1\", \"content\": \"- Item 1\\n- Item 2\\n- Item 3\"}, {\"title\": \"Slide 2\", \"content\": \"- Item A\\n- Item B\"}]}"],
    image_filename: Annotated[str, "Optional image filename to integrate into the presentation"] = "",
    image_data_base64: Annotated[str, "Framework may supply Base64 image content as fallback"] = ""
) -> Dict[str, Any]:
    """
    Converts JSON input to PowerPoint presentation with support for bullet point lists and optional image integration
    
    Args:
        input_data: JSON string containing slide data in this format:
            {"slides": [
                {"title": "Slide 1", "content": "- Item 1\\n- Item 2\\n- Item 3"},
                {"title": "Slide 2", "content": "- Item A\\n- Item B"}
            ]}
        image_filename: Optional image filename to integrate into the presentation
        image_data_base64: Framework may supply Base64 image content as fallback
    
    Returns:
        Dictionary with 'results' and 'artifacts' keys:
        - 'results': Success message or error message
        - 'artifacts': List of artifact dictionaries with 'name', 'b64', and 'mime' keys
    """
    print("Starting json_to_pptx execution...")
    try:
        import json
        data = json.loads(input_data)
        
        if not isinstance(data, dict) or 'slides' not in data:
            return {"results": {"error": "Input must be a JSON object containing 'slides' array"}}
            
        slides = data['slides']
        print(f"Processing {len(slides)} slides...")
        
        # Load image if provided
        image_bytes = None
        if image_filename:
            image_bytes = _load_image_bytes(image_filename, image_data_base64)
            if image_bytes:
                print(f"Loaded image: {image_filename}")
            else:
                print(f"Failed to load image: {image_filename}")
        
        # Create presentation
        prs = Presentation()
        print("Created PowerPoint presentation object")
        
        for i, slide in enumerate(slides):
            title = slide.get('title', 'Untitled Slide')
            content = slide.get('content', '')
            
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide_obj = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide_obj.shapes.title
            title_shape.text = title
            print(f"Added slide {i+1}: {title}")
            
            # Add content
            body_shape = slide_obj.placeholders[1]
            tf = body_shape.text_frame
            tf.text = ""
            
            # Set text alignment to left
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            
            # Process bullet points if content contains them
            if content.strip() and content.strip().startswith('-'):
                # Split by newline and process each line
                items = [item.strip() for item in content.split('\n') if item.strip()]
                print(f"Slide {i+1} has {len(items)} bullet points")
                for item in items:
                    if item.startswith('-'):
                        item_text = item[1:].strip()  # Remove the dash
                        p = tf.add_paragraph()
                        p.text = item_text
                        p.level = 0
                        p.font.size = Pt(24)
                        p.space_after = Pt(6)
                        p.alignment = PP_ALIGN.LEFT
            else:
                # Handle regular text without bullet points
                p = tf.add_paragraph()
                p.text = content
                p.font.size = Pt(24)
                p.space_after = Pt(6)
                p.alignment = PP_ALIGN.LEFT
            
            # Add image to first slide if provided
            if i == 0 and image_bytes:
                _add_image_to_slide(slide_obj, image_bytes, 
                                   left=Inches(0.5), top=Inches(3.5), 
                                   width=Inches(4), height=Inches(3))
            
        # Save presentation
        pptx_output_path = os.path.join(os.getcwd(), "output_presentation.pptx")
        prs.save(pptx_output_path)
        print(f"Saved PowerPoint presentation to {pptx_output_path}")
        
        # Convert to PDF using pdfkit
        pdf_output_path = os.path.join(os.getcwd(), "output_presentation.pdf")
        print(f"Starting PDF conversion to {pdf_output_path}")
        
        # Create HTML representation of the presentation
        html_content = "<html><head><title>PowerPoint Presentation</title></head><body>"
        
        for i, slide in enumerate(slides):
            title = slide.get('title', 'Untitled Slide')
            content = slide.get('content', '')
            
            html_content += f"<div style='margin-bottom: 20px; border: 1px solid #ddd; padding: 10px; border-radius: 4px;'><h2 style='color: #2c3e50; margin-bottom: 10px;'>{title}</h2>"
            
            if content.strip() and content.strip().startswith('-'):
                items = [item.strip() for item in content.split('\n') if item.strip()]
                html_content += "<ul style='margin: 0; padding-left: 20px;'>"
                for item in items:
                    if item.startswith('-'):
                        item_text = item[1:].strip()
                        html_content += f"<li style='margin-bottom: 5px; font-size: 18px;'>{item_text}</li>"
                html_content += "</ul>"
            else:
                html_content += f"<p style='margin: 0; font-size: 18px;'>{content}</p>"
            
            html_content += "</div>"
        
        html_content += "</body></html>"
        
        # Generate PDF from HTML
        try:
            print("Converting HTML to PDF...")
            from_string(html_content, pdf_output_path, options={"page-size": "A4", "margin-top": "0.75in", "margin-right": "0.75in", "margin-bottom": "0.75in", "margin-left": "0.75in"})
            print(f"PDF successfully created at {pdf_output_path}")
        except Exception as e:
            # If pdfkit fails, return only PPTX
            print(f"Warning: PDF conversion failed: {str(e)}")
            # Remove the PDF file if it exists
            if os.path.exists(pdf_output_path):
                os.remove(pdf_output_path)
            print("PDF file removed due to conversion error")
        
        # Read PPTX file as bytes
        with open(pptx_output_path, "rb") as f:
            pptx_bytes = f.read()
            
        # Encode PPTX as base64
        pptx_b64 = base64.b64encode(pptx_bytes).decode('utf-8')
        print("PPTX file successfully encoded to base64")
        
        # Read PDF file as bytes if it exists
        pdf_b64 = None
        if os.path.exists(pdf_output_path):
            with open(pdf_output_path, "rb") as f:
                pdf_bytes = f.read()
            pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
            print("PDF file successfully encoded to base64")
        
        # Prepare artifacts
        artifacts = [
            {
                "name": "presentation.pptx",
                "b64": pptx_b64,
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            }
        ]
        
        # Add PDF if conversion was successful
        if pdf_b64:
            artifacts.append({
                "name": "presentation.pdf",
                "b64": pdf_b64,
                "mime": "application/pdf",
            })
            print(f"Added {len(artifacts)} artifacts to response")
        else:
            print("No PDF artifact added due to conversion failure")
        
        return {
            "results": {
                "operation": "json_to_pptx",
                "message": "PowerPoint presentation generated successfully.",
                "pdf_generated": pdf_b64 is not None,
                "image_included": image_bytes is not None,
            },
            "artifacts": artifacts,
            "display": {
                "open_canvas": True,
                "primary_file": "presentation.pptx",
                "mode": "replace",
                "viewer_hint": "powerpoint",
            },
            "meta_data": {
                "generated_slides": len(slides),
                "output_files": [f"presentation.pptx", "presentation.pdf"] if pdf_b64 else ["presentation.pptx"],
                "output_file_paths": [pptx_output_path, pdf_output_path] if pdf_b64 else [pptx_output_path],
            },
        }
    except Exception as e:
        print(f"Error in json_to_pptx: {str(e)}")
        return {"results": {"error": f"Error creating PowerPoint: {str(e)}"}}


@mcp.tool
def markdown_to_pptx(
    markdown_content: Annotated[str, "Markdown content with headers (# or ##) as slide titles and content below each header"],
    image_filename: Annotated[str, "Optional image filename to integrate into the presentation"] = "",
    image_data_base64: Annotated[str, "Framework may supply Base64 image content as fallback"] = ""
) -> Dict[str, Any]:
    """
    Converts markdown content to PowerPoint presentation with support for bullet point lists and optional image integration
    
    Args:
        markdown_content: Markdown content where headers (# or ##) become slide titles and content below becomes slide content
        image_filename: Optional image filename to integrate into the presentation
        image_data_base64: Framework may supply Base64 image content as fallback
    
    Returns:
        Dictionary with 'results' and 'artifacts' keys:
        - 'results': Success message or error message
        - 'artifacts': List of artifact dictionaries with 'name', 'b64', and 'mime' keys
    """
    print("Starting markdown_to_pptx execution...")
    try:
        # Parse markdown into slides
        slides = _parse_markdown_slides(markdown_content)
        print(f"Parsed {len(slides)} slides from markdown")
        
        if not slides:
            return {"results": {"error": "No slides could be parsed from markdown content"}}
        
        # Load image if provided
        image_bytes = None
        if image_filename:
            image_bytes = _load_image_bytes(image_filename, image_data_base64)
            if image_bytes:
                print(f"Loaded image: {image_filename}")
            else:
                print(f"Failed to load image: {image_filename}")
        
        # Create presentation
        prs = Presentation()
        print("Created PowerPoint presentation object")
        
        for i, slide_data in enumerate(slides):
            title = slide_data.get('title', 'Untitled Slide')
            content = slide_data.get('content', '')
            
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide_obj = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide_obj.shapes.title
            title_shape.text = title
            print(f"Added slide {i+1}: {title}")
            
            # Add content
            body_shape = slide_obj.placeholders[1]
            tf = body_shape.text_frame
            tf.text = ""
            
            # Set text alignment to left
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            
            # Process content - handle bullet points and regular text
            if content.strip():
                lines = content.split('\n')
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    p = tf.add_paragraph()
                    
                    # Handle bullet points (- or *)
                    if line.startswith(('- ', '* ')):
                        p.text = line[2:].strip()
                        p.level = 0
                    # Handle sub-bullets (indented)
                    elif line.startswith(('  - ', '  * ', '\t- ', '\t* ')):
                        p.text = line.strip()[2:].strip()
                        p.level = 1
                    else:
                        # Regular text
                        p.text = line
                        p.level = 0
                    
                    p.font.size = Pt(24)
                    p.space_after = Pt(6)
                    p.alignment = PP_ALIGN.LEFT
            
            # Add image to first slide if provided
            if i == 0 and image_bytes:
                _add_image_to_slide(slide_obj, image_bytes, 
                                   left=Inches(0.5), top=Inches(3.5), 
                                   width=Inches(4), height=Inches(3))
        
        # Save presentation
        pptx_output_path = os.path.join(os.getcwd(), "output_presentation.pptx")
        prs.save(pptx_output_path)
        print(f"Saved PowerPoint presentation to {pptx_output_path}")
        
        # Convert to PDF using pdfkit
        pdf_output_path = os.path.join(os.getcwd(), "output_presentation.pdf")
        print(f"Starting PDF conversion to {pdf_output_path}")
        
        # Create HTML representation of the presentation
        html_content = "<html><head><title>PowerPoint Presentation</title></head><body>"
        
        for i, slide_data in enumerate(slides):
            title = slide_data.get('title', 'Untitled Slide')
            content = slide_data.get('content', '')
            
            html_content += f"<div style='margin-bottom: 20px; border: 1px solid #ddd; padding: 10px; border-radius: 4px;'><h2 style='color: #2c3e50; margin-bottom: 10px;'>{title}</h2>"
            
            if content.strip():
                lines = content.split('\n')
                bullet_lines = []
                regular_lines = []
                
                for line in lines:
                    line = line.strip()
                    if line.startswith(('- ', '* ')):
                        bullet_lines.append(line[2:].strip())
                    elif line.startswith(('  - ', '  * ', '\t- ', '\t* ')):
                        bullet_lines.append(f"&nbsp;&nbsp;&bull; {line.strip()[2:].strip()}")
                    elif line:
                        regular_lines.append(line)
                
                if bullet_lines:
                    html_content += "<ul style='margin: 0; padding-left: 20px;'>"
                    for item in bullet_lines:
                        html_content += f"<li style='margin-bottom: 5px; font-size: 18px;'>{item}</li>"
                    html_content += "</ul>"
                
                if regular_lines:
                    for line in regular_lines:
                        html_content += f"<p style='margin: 5px 0; font-size: 18px;'>{line}</p>"
            
            html_content += "</div>"
        
        html_content += "</body></html>"
        
        # Generate PDF from HTML
        try:
            print("Converting HTML to PDF...")
            from_string(html_content, pdf_output_path, options={"page-size": "A4", "margin-top": "0.75in", "margin-right": "0.75in", "margin-bottom": "0.75in", "margin-left": "0.75in"})
            print(f"PDF successfully created at {pdf_output_path}")
        except Exception as e:
            # If pdfkit fails, return only PPTX
            print(f"Warning: PDF conversion failed: {str(e)}")
            # Remove the PDF file if it exists
            if os.path.exists(pdf_output_path):
                os.remove(pdf_output_path)
            print("PDF file removed due to conversion error")
        
        # Read PPTX file as bytes
        with open(pptx_output_path, "rb") as f:
            pptx_bytes = f.read()
            
        # Encode PPTX as base64
        pptx_b64 = base64.b64encode(pptx_bytes).decode('utf-8')
        print("PPTX file successfully encoded to base64")
        
        # Read PDF file as bytes if it exists
        pdf_b64 = None
        if os.path.exists(pdf_output_path):
            with open(pdf_output_path, "rb") as f:
                pdf_bytes = f.read()
            pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
            print("PDF file successfully encoded to base64")
        
        # Prepare artifacts
        artifacts = [
            {
                "name": "presentation.pptx",
                "b64": pptx_b64,
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            }
        ]
        
        # Add PDF if conversion was successful
        if pdf_b64:
            artifacts.append({
                "name": "presentation.pdf",
                "b64": pdf_b64,
                "mime": "application/pdf",
            })
            print(f"Added {len(artifacts)} artifacts to response")
        else:
            print("No PDF artifact added due to conversion failure")
        
        return {
            "results": {
                "operation": "markdown_to_pptx",
                "message": "PowerPoint presentation generated successfully from markdown.",
                "pdf_generated": pdf_b64 is not None,
                "image_included": image_bytes is not None,
            },
            "artifacts": artifacts,
            "display": {
                "open_canvas": True,
                "primary_file": "presentation.pptx",
                "mode": "replace",
                "viewer_hint": "powerpoint",
            },
            "meta_data": {
                "generated_slides": len(slides),
                "output_files": [f"presentation.pptx", "presentation.pdf"] if pdf_b64 else ["presentation.pptx"],
                "output_file_paths": [pptx_output_path, pdf_output_path] if pdf_b64 else [pptx_output_path],
            },
        }
    except Exception as e:
        print(f"Error in markdown_to_pptx: {str(e)}")
        return {"results": {"error": f"Error creating PowerPoint from markdown: {str(e)}"}}


if __name__ == "__main__":
    mcp.run()
